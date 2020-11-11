from pptx import Presentation
from pptx.util import Inches,Pt
import os
import tkinter as tk

def makePPT():
    week=name.get()
    prs = Presentation("template.pptx")
    pleft=Inches(1)
    ptop=Inches(1.6)
    pwidth=Inches(7.5)
    pheight=Inches(5.5)
    j=1
    for i in os.listdir('images'):
        full_path = os.path.join('images\\', i)
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title=slide.shapes.title
        title.text='예제 '+ os.path.splitext(i)[0]
        slide_num=slide.placeholders[1]
        slide_num.text=str(j)
        slide.shapes.add_picture(full_path,pleft,ptop,pwidth,pheight)
        j+=1
        filename=str(week)+'주차 김주연.pptx'
        prs.save(filename)
        return exit

root = tk.Tk()
root.geometry("200x80")
root.title("make ppt automation")
textLable=tk.Label(root,text="input week : ")
textLable.grid(column=0,row=0)

name = tk.StringVar()
nameEntered = tk.Entry(root, width = 15, textvariable = name)
nameEntered.grid(column = 1, row = 0)

btnRead=tk.Button(root, height=1, width=10, text="create", command=makePPT)
btnRead.grid(columnspan=2,row=1)
root.mainloop()



